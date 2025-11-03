# app.py
import uvicorn
from fastapi import FastAPI, UploadFile, Form
from fastapi.responses import FileResponse
from pydantic import BaseModel
from transformers import pipeline
from pptx import Presentation
from tempfile import NamedTemporaryFile
import os
import gradio as gr
from tqdm import tqdm

# Load both translation directions
model_ko_en = "Helsinki-NLP/opus-mt-ko-en"
model_en_ko = "seongs/ke-t5-base-aihub-koen-translation-integrated-10m-en-to-ko"

pipe_ko_en = pipeline(
    "translation", model=model_ko_en, tokenizer=model_ko_en, max_length=512
)
pipe_en_ko = pipeline(
    "translation", model=model_en_ko, tokenizer=model_en_ko, max_length=512
)

app = FastAPI(title="Korean-English Translator API + PPT Tool")

# ------------------- API TRANSLATION ENDPOINT -------------------


class TranslateRequest(BaseModel):
    text: str
    src_lang: str  # "ko" or "en"


@app.post("/translate")
def translate(req: TranslateRequest):
    if req.src_lang == "ko":
        out = pipe_ko_en(req.text)[0]["translation_text"]
    elif req.src_lang == "en":
        out = pipe_en_ko(req.text)[0]["translation_text"]
    else:
        return {"error": "unsupported source language"}
    return {"translation": out}


# ------------------- PPT TRANSLATION HANDLER -------------------


def translate_text(text, src_lang):
    if not text.strip():
        return text
    if src_lang == "ko":
        return pipe_ko_en(text)[0]["translation_text"]
    else:
        return pipe_en_ko(text)[0]["translation_text"]


def translate_pptx_file(input_path, src_lang="en"):
    prs = Presentation(input_path)
    for slide in tqdm(prs.slides, desc="Translating slides"):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.text = translate_text(run.text, src_lang)
    # Save to temporary file
    tmp_out = NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp_out.name)
    return tmp_out.name


@app.post("/translate_pptx")
async def translate_pptx_endpoint(file: UploadFile, src_lang: str = Form("en")):
    tmp_input = NamedTemporaryFile(delete=False, suffix=".pptx")
    tmp_input.write(await file.read())
    tmp_input.close()

    output_path = translate_pptx_file(tmp_input.name, src_lang)
    translated_name = f"translated_{os.path.basename(file.filename)}"
    return FileResponse(
        output_path,
        filename=translated_name,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )


# ------------------- GRADIO WEB UI -------------------


def gradio_translate_text(text, src):
    return translate_text(text, src)


def gradio_translate_ppt(file, src):
    if file is None:
        return "Please upload a PPTX file.", None
    out_path = translate_pptx_file(file.name, src)
    return f"✅ Translation complete. Download below.", out_path


with gr.Blocks(title="Korean-English Translator") as demo:
    gr.Markdown("## 🈳 Korean ↔ English Translator\nSupports text and PPTX files.")
    with gr.Tab("Text Translation"):
        input_text = gr.Textbox(label="Input Text")
        lang = gr.Radio(["en", "ko"], label="Source Language", value="en")
        output_text = gr.Textbox(label="Translation")
        translate_btn = gr.Button("Translate")
        translate_btn.click(
            fn=gradio_translate_text, inputs=[input_text, lang], outputs=output_text
        )

    with gr.Tab("PPTX Translation"):
        ppt_file = gr.File(label="Upload PowerPoint File (.pptx)", file_types=[".pptx"])
        ppt_lang = gr.Radio(["en", "ko"], label="Source Language", value="en")
        ppt_output = gr.File(label="Download Translated PPTX")
        ppt_status = gr.Textbox(label="Status")
        ppt_btn = gr.Button("Translate PPTX")
        ppt_btn.click(
            fn=gradio_translate_ppt,
            inputs=[ppt_file, ppt_lang],
            outputs=[ppt_status, ppt_output],
        )

demo.queue().launch(server_name="0.0.0.0", server_port=7860)

if __name__ == "__main__":
    uvicorn.run(app, host="0.0.0.0", port=8000)
